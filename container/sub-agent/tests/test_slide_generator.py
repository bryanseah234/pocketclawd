"""Tests for slide generator module."""
import pytest
import json


@pytest.mark.skipif(
    not __import__("importlib").util.find_spec("pptx"),
    reason="python-pptx not installed"
)
class TestSlideGenerator:
    def test_slide_template_enum_has_all_4(self):
        from src.slide_generator import SlideTemplate
        names = {t.value for t in SlideTemplate}
        assert "executive_summary" in names
        assert "technical_deep_dive" in names
        assert "project_proposal" in names
        assert "data_report" in names

    def test_generate_returns_bytes_for_each_template(self):
        from src.slide_generator import SlideTemplate, PresentationRequest, SlideContent, generate_presentation
        for template in SlideTemplate:
            req = PresentationRequest(
                title="Test",
                template=template,
                slides=[SlideContent(title="Slide 1", bullet_points=["Point A"])]
            )
            result = generate_presentation(req)
            assert isinstance(result, bytes)
            assert len(result) > 100

    def test_generate_has_correct_slide_count(self):
        from pptx import Presentation
        from src.slide_generator import SlideTemplate, PresentationRequest, SlideContent, generate_presentation
        import io
        req = PresentationRequest(
            title="Test", template=SlideTemplate.DATA_REPORT,
            slides=[SlideContent(title="S1"), SlideContent(title="S2")]
        )
        buf = generate_presentation(req)
        prs = Presentation(io.BytesIO(buf))
        assert len(prs.slides) == 3  # 1 title + 2 content

    def test_generate_with_empty_slides_returns_title_only(self):
        from pptx import Presentation
        from src.slide_generator import SlideTemplate, PresentationRequest, generate_presentation
        import io
        req = PresentationRequest(title="Empty", template=SlideTemplate.EXECUTIVE_SUMMARY, slides=[])
        buf = generate_presentation(req)
        prs = Presentation(io.BytesIO(buf))
        assert len(prs.slides) == 1

    def test_generated_pptx_is_valid(self):
        from pptx import Presentation
        from src.slide_generator import SlideTemplate, PresentationRequest, SlideContent, generate_presentation
        import io
        req = PresentationRequest(
            title="Valid", template=SlideTemplate.PROJECT_PROPOSAL,
            slides=[SlideContent(title="Intro", bullet_points=["A", "B"])]
        )
        buf = generate_presentation(req)
        prs = Presentation(io.BytesIO(buf))
        assert prs is not None

    def test_parse_slide_request_returns_request_from_valid_json(self):
        from src.slide_generator import parse_slide_request, SlideTemplate
        llm_json = json.dumps({
            "title": "Q1 Review",
            "template": "data_report",
            "slides": [{"title": "Revenue", "bullet_points": ["Up 10%"]}]
        })
        result = parse_slide_request("make a slide", llm_json)
        assert result is not None
        assert result.title == "Q1 Review"
        assert result.template == SlideTemplate.DATA_REPORT

    def test_parse_slide_request_returns_none_for_invalid_json(self):
        from src.slide_generator import parse_slide_request
        result = parse_slide_request("make a slide", "not json {{{{")
        assert result is None

    def test_parse_slide_request_returns_none_for_missing_title(self):
        from src.slide_generator import parse_slide_request
        llm_json = json.dumps({"template": "data_report", "slides": []})
        result = parse_slide_request("make a slide", llm_json)
        assert result is None

    def test_parse_slide_request_accepts_unknown_template_with_default(self):
        from src.slide_generator import parse_slide_request, SlideTemplate
        llm_json = json.dumps({"title": "T", "template": "unknown_template", "slides": []})
        result = parse_slide_request("t", llm_json)
        assert result is not None
        assert result.template == SlideTemplate.EXECUTIVE_SUMMARY

    def test_generate_includes_title_slide(self):
        from pptx import Presentation
        from src.slide_generator import SlideTemplate, PresentationRequest, generate_presentation
        import io
        req = PresentationRequest(title="My Title", template=SlideTemplate.TECHNICAL_DEEP_DIVE, slides=[])
        buf = generate_presentation(req)
        prs = Presentation(io.BytesIO(buf))
        assert len(prs.slides) >= 1
