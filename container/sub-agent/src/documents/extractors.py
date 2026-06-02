"""
Text extraction functions for various document formats.

Supports PDF (text + OCR fallback), DOCX, CSV, TXT, and images (OCR).

Requirements: REQ-5.1
"""

import csv
import io
import logging
from typing import Callable

# Module-level imports so unittest.mock.patch can target these names.
# Wrapped in try/except so missing optional deps do not break module import.
try:  # pragma: no cover
    from PyPDF2 import PdfReader  # type: ignore
except Exception:  # pragma: no cover
    PdfReader = None  # type: ignore
try:  # pragma: no cover
    from docx import Document  # type: ignore
except Exception:  # pragma: no cover
    Document = None  # type: ignore
try:  # pragma: no cover
    from PIL import Image  # type: ignore
except Exception:  # pragma: no cover
    Image = None  # type: ignore
try:  # pragma: no cover
    import pytesseract  # type: ignore
except Exception:  # pragma: no cover
    pytesseract = None  # type: ignore
try:  # pragma: no cover
    from openpyxl import load_workbook  # type: ignore
except Exception:  # pragma: no cover
    load_workbook = None  # type: ignore
try:  # pragma: no cover
    from pptx import Presentation  # type: ignore
except Exception:  # pragma: no cover
    Presentation = None  # type: ignore
try:  # pragma: no cover
    from bs4 import BeautifulSoup  # type: ignore
except Exception:  # pragma: no cover
    BeautifulSoup = None  # type: ignore
# Register HEIC/HEIF support with Pillow so iPhone photos open transparently.
try:  # pragma: no cover
    import pillow_heif  # type: ignore
    pillow_heif.register_heif_opener()
except Exception:  # pragma: no cover
    pass

logger = logging.getLogger(__name__)

# Content type to extractor mapping
_EXTRACTOR_MAP: dict[str, Callable[[bytes], str]] = {}


def extract_pdf(content: bytes) -> str:
    """
    Extract text from a PDF file.

    Uses PyPDF2 for text-based PDFs. Falls back to pytesseract OCR
    for scanned pages that yield no text.

    Args:
        content: Raw PDF file bytes.

    Returns:
        Extracted text content.
    """
    # PdfReader is imported at module level for patchability

    reader = PdfReader(io.BytesIO(content))
    pages_text: list[str] = []

    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            pages_text.append(text)
        else:
            # Scanned page — attempt OCR
            ocr_text = _ocr_pdf_page(content, reader.pages.index(page))
            if ocr_text:
                pages_text.append(ocr_text)

    return "\n\n".join(pages_text)


def _ocr_pdf_page(content: bytes, page_index: int) -> str:
    """
    OCR a single PDF page using pytesseract.

    Converts the page to an image and runs OCR. Returns empty string
    if OCR dependencies are unavailable.
    """
    try:
        from pdf2image import convert_from_bytes
        # pytesseract is imported at module level for patchability

        images = convert_from_bytes(content, first_page=page_index + 1, last_page=page_index + 1)
        if images:
            return pytesseract.image_to_string(images[0])
    except ImportError:
        logger.warning("pdf2image or pytesseract not available for OCR fallback")
    except Exception as e:
        logger.warning("OCR failed for page %d: %s", page_index, e)

    return ""


def extract_docx(content: bytes) -> str:
    """
    Extract text from a DOCX file.

    Uses python-docx to read all paragraphs.

    Args:
        content: Raw DOCX file bytes.

    Returns:
        Extracted text content with paragraphs separated by newlines.
    """
    # Document is imported at module level for patchability

    doc = Document(io.BytesIO(content))
    paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
    return "\n\n".join(paragraphs)


def extract_csv(content: bytes) -> str:
    """
    Extract text from a CSV file.

    Reads the CSV and converts each row to a readable text line.

    Args:
        content: Raw CSV file bytes.

    Returns:
        Text representation of the CSV data.
    """
    text_content = content.decode("utf-8")
    reader = csv.reader(io.StringIO(text_content))

    rows: list[str] = []
    headers: list[str] | None = None

    for i, row in enumerate(reader):
        if i == 0:
            headers = row
            rows.append(" | ".join(row))
        else:
            if headers:
                # Format as "header: value" pairs for better context
                pairs = [f"{h}: {v}" for h, v in zip(headers, row) if v.strip()]
                rows.append(", ".join(pairs))
            else:
                rows.append(" | ".join(row))

    return "\n".join(rows)


def extract_txt(content: bytes) -> str:
    """
    Extract text from a plain text file.

    Decodes as UTF-8.

    Args:
        content: Raw text file bytes.

    Returns:
        Decoded text content.
    """
    return content.decode("utf-8")


def extract_image(content: bytes) -> str:
    """
    Extract text + description from an image.

    F1 (Wave 9): Routes through `vision.bedrock_vision.describe_image` which
    honours the CLAWD_VISION_PROVIDER env var:
    - "bedrock" (default in cloud) — Claude Sonnet 4.5 native vision
    - "tesseract" — legacy pytesseract OCR
    - "auto" — try Bedrock; fall back to tesseract on any failure

    The Bedrock provider returns "TEXT:\n...\n\nDESCRIPTION:\n..." which
    is preserved verbatim so embedders/index get both signals.

    Args:
        content: Raw image file bytes (PNG, JPEG, etc.).

    Returns:
        Extracted text + description.
    """
    # Late import so the legacy tesseract-only test path keeps working.
    from src.vision.bedrock_vision import describe_image

    return describe_image(content)


def extract_xlsx(content: bytes) -> str:
    """Extract text from an XLSX workbook (all sheets, header-aware rows)."""
    if load_workbook is None:
        raise ValueError("openpyxl not available for XLSX extraction")
    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    out: list[str] = []
    for ws in wb.worksheets:
        out.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                out.append(" | ".join(cells))
    try:
        wb.close()
    except Exception:
        pass
    return "\n".join(out)


def extract_pptx(content: bytes) -> str:
    """Extract text from a PPTX deck (all slides, all text frames + notes)."""
    if Presentation is None:
        raise ValueError("python-pptx not available for PPTX extraction")
    prs = Presentation(io.BytesIO(content))
    out: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        out.append(f"# Slide {idx}")
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs).strip()
                    if txt:
                        out.append(txt)
        notes = getattr(slide, "notes_slide", None)
        if notes is not None and getattr(notes, "notes_text_frame", None) is not None:
            note_txt = notes.notes_text_frame.text.strip()
            if note_txt:
                out.append(f"(notes) {note_txt}")
    return "\n".join(out)


def extract_html(content: bytes) -> str:
    """Extract readable text from an HTML file, stripping scripts/styles."""
    text = content.decode("utf-8", errors="replace")
    if BeautifulSoup is None:
        # Fallback: crude tag strip
        import re
        return re.sub(r"<[^>]+>", " ", text)
    soup = BeautifulSoup(text, "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def extract_text(content: bytes, content_type: str) -> str:
    """
    Route text extraction to the appropriate extractor based on content type.

    Supported content types:
    - application/pdf → extract_pdf
    - application/vnd.openxmlformats-officedocument.wordprocessingml.document → extract_docx
    - text/csv → extract_csv
    - text/plain → extract_txt
    - image/* → extract_image

    Args:
        content: Raw file bytes.
        content_type: MIME type of the file.

    Returns:
        Extracted text content.

    Raises:
        ValueError: If the content type is not supported.
    """
    extractor = _get_extractor(content_type)
    if extractor is None:
        raise ValueError(f"Unsupported content type: {content_type}")
    return extractor(content)


def _get_extractor(content_type: str) -> Callable[[bytes], str] | None:
    """Look up the extractor function for a given content type."""
    # Normalize content type (strip parameters like charset)
    ct = content_type.split(";")[0].strip().lower()

    mapping: dict[str, Callable[[bytes], str]] = {
        "application/pdf": extract_pdf,
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": extract_docx,
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": extract_xlsx,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": extract_pptx,
        "text/csv": extract_csv,
        "text/plain": extract_txt,
        "text/markdown": extract_txt,
        "text/html": extract_html,
        "application/xhtml+xml": extract_html,
        "image/png": extract_image,
        "image/jpeg": extract_image,
        "image/jpg": extract_image,
        "image/tiff": extract_image,
        "image/bmp": extract_image,
        "image/gif": extract_image,
        "image/webp": extract_image,
        "image/heic": extract_image,
        "image/heif": extract_image,
    }

    return mapping.get(ct)


def is_supported(content_type: str) -> bool:
    """True if a usable extractor exists for this MIME type.

    Used by the indexer to fail gracefully (friendly user message) instead of
    raising ValueError on an unsupported upload.
    """
    return _get_extractor(content_type) is not None
