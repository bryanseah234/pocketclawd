"""
PPTX slide generation from document summaries.

Supports 4 templates:
  - Corporate: Blue/gray, formal fonts, company logo placeholder
  - Modern: Clean white, minimal, sans-serif
  - Elegant: Dark background, gold accents, serif fonts
  - Informative: Green/white, data-focused, charts-friendly

The generator:
  1. Accepts a summary text and template choice
  2. Parses the summary into slide sections (title + bullet points)
  3. Generates a PPTX file using python-pptx
  4. Returns the file bytes for upload to S3

Requirements: REQ-5.3
"""

import io
import re
from dataclasses import dataclass, field
from enum import Enum
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


# ---------------------------------------------------------------------------
# Types
# ---------------------------------------------------------------------------


class SlideTemplate(str, Enum):
    """Available slide templates."""

    CORPORATE = "corporate"
    MODERN = "modern"
    ELEGANT = "elegant"
    INFORMATIVE = "informative"


@dataclass
class SlideSection:
    """A parsed section of the summary representing one slide."""

    title: str
    bullets: list[str] = field(default_factory=list)


@dataclass
class SlideResult:
    """Result of slide generation."""

    file_bytes: bytes
    filename: str
    slide_count: int
    template: SlideTemplate


# ---------------------------------------------------------------------------
# Template styling definitions
# ---------------------------------------------------------------------------


@dataclass
class TemplateStyle:
    """Visual styling for a slide template."""

    # Background
    bg_color: RGBColor
    # Title
    title_font: str
    title_size: Pt
    title_color: RGBColor
    title_bold: bool
    # Body
    body_font: str
    body_size: Pt
    body_color: RGBColor
    # Accent
    accent_color: RGBColor
    # Footer
    footer_text: str


TEMPLATE_STYLES: dict[SlideTemplate, TemplateStyle] = {
    SlideTemplate.CORPORATE: TemplateStyle(
        bg_color=RGBColor(0xF0, 0xF4, 0xF8),  # Light blue-gray
        title_font="Calibri",
        title_size=Pt(32),
        title_color=RGBColor(0x1A, 0x36, 0x5D),  # Dark navy
        title_bold=True,
        body_font="Calibri",
        body_size=Pt(18),
        body_color=RGBColor(0x2C, 0x3E, 0x50),  # Dark gray-blue
        accent_color=RGBColor(0x2E, 0x86, 0xC1),  # Corporate blue
        footer_text="Confidential — Company Presentation",
    ),
    SlideTemplate.MODERN: TemplateStyle(
        bg_color=RGBColor(0xFF, 0xFF, 0xFF),  # Clean white
        title_font="Arial",
        title_size=Pt(36),
        title_color=RGBColor(0x21, 0x21, 0x21),  # Near-black
        title_bold=False,
        body_font="Arial",
        body_size=Pt(18),
        body_color=RGBColor(0x42, 0x42, 0x42),  # Dark gray
        accent_color=RGBColor(0x00, 0xBC, 0xD4),  # Teal accent
        footer_text="",
    ),
    SlideTemplate.ELEGANT: TemplateStyle(
        bg_color=RGBColor(0x1C, 0x1C, 0x2E),  # Dark background
        title_font="Georgia",
        title_size=Pt(34),
        title_color=RGBColor(0xD4, 0xAF, 0x37),  # Gold
        title_bold=True,
        body_font="Georgia",
        body_size=Pt(18),
        body_color=RGBColor(0xEC, 0xEC, 0xEC),  # Light gray
        accent_color=RGBColor(0xD4, 0xAF, 0x37),  # Gold accent
        footer_text="",
    ),
    SlideTemplate.INFORMATIVE: TemplateStyle(
        bg_color=RGBColor(0xF9, 0xFF, 0xF9),  # Light green-white
        title_font="Verdana",
        title_size=Pt(30),
        title_color=RGBColor(0x1B, 0x5E, 0x20),  # Dark green
        title_bold=True,
        body_font="Verdana",
        body_size=Pt(16),
        body_color=RGBColor(0x2E, 0x7D, 0x32),  # Medium green
        accent_color=RGBColor(0x4C, 0xAF, 0x50),  # Green accent
        footer_text="Data-Driven Insights",
    ),
}


# ---------------------------------------------------------------------------
# Summary parser
# ---------------------------------------------------------------------------


def parse_summary_to_sections(summary: str) -> list[SlideSection]:
    """
    Parse a summary text into slide sections.

    Supports multiple formats:
      - Markdown-style headers (# or ##) followed by bullet points (- or *)
      - Numbered sections (1. Title) followed by bullet points
      - Plain paragraphs separated by double newlines (each becomes a slide)

    Returns at least one section (the title slide) even for empty input.
    """
    if not summary or not summary.strip():
        return [SlideSection(title="Presentation", bullets=[])]

    lines = summary.strip().split("\n")
    sections: list[SlideSection] = []
    current_section: Optional[SlideSection] = None

    for line in lines:
        stripped = line.strip()

        # Skip empty lines
        if not stripped:
            continue

        # Check for markdown headers: # Title or ## Title
        header_match = re.match(r"^#{1,3}\s+(.+)$", stripped)
        if header_match:
            if current_section is not None:
                sections.append(current_section)
            current_section = SlideSection(title=header_match.group(1).strip())
            continue

        # Check for numbered sections: 1. Title or 1) Title
        numbered_match = re.match(r"^\d+[.)]\s+(.+)$", stripped)
        if numbered_match and not stripped.startswith("  "):
            # Only treat as a section header if it's not indented (not a sub-bullet)
            if current_section is not None:
                sections.append(current_section)
            current_section = SlideSection(title=numbered_match.group(1).strip())
            continue

        # Check for bullet points: - item or * item or • item
        bullet_match = re.match(r"^[-*•]\s+(.+)$", stripped)
        if bullet_match:
            if current_section is None:
                current_section = SlideSection(title="Overview")
            current_section.bullets.append(bullet_match.group(1).strip())
            continue

        # Check for indented bullets (sub-items under numbered sections)
        indented_bullet = re.match(r"^\s+[-*•]\s+(.+)$", line)
        if indented_bullet:
            if current_section is None:
                current_section = SlideSection(title="Overview")
            current_section.bullets.append(indented_bullet.group(1).strip())
            continue

        # Plain text — treat as a new section title if no current section,
        # otherwise add as a bullet point
        if current_section is None:
            current_section = SlideSection(title=stripped)
        else:
            current_section.bullets.append(stripped)

    # Don't forget the last section
    if current_section is not None:
        sections.append(current_section)

    # Ensure at least one section
    if not sections:
        sections = [SlideSection(title="Presentation", bullets=[])]

    return sections


# ---------------------------------------------------------------------------
# PPTX generation
# ---------------------------------------------------------------------------


class SlideGenerator:
    """Generates PPTX presentations from parsed summary sections."""

    def __init__(self, template: SlideTemplate = SlideTemplate.MODERN):
        self.template = template
        self.style = TEMPLATE_STYLES[template]

    def generate(self, summary: str, title: Optional[str] = None) -> SlideResult:
        """
        Generate a PPTX presentation from a summary text.

        Args:
            summary: The summary text to convert into slides.
            title: Optional presentation title. If not provided, uses the
                   first section title.

        Returns:
            SlideResult with the PPTX file bytes and metadata.
        """
        sections = parse_summary_to_sections(summary)

        # Use provided title or first section title
        presentation_title = title or sections[0].title

        prs = Presentation()

        # Set slide dimensions (widescreen 16:9)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Generate title slide
        self._add_title_slide(prs, presentation_title)

        # Generate content slides
        for section in sections:
            self._add_content_slide(prs, section)

        # Save to bytes
        buffer = io.BytesIO()
        prs.save(buffer)
        file_bytes = buffer.getvalue()

        filename = self._sanitize_filename(presentation_title)

        return SlideResult(
            file_bytes=file_bytes,
            filename=filename,
            slide_count=len(sections) + 1,  # +1 for title slide
            template=self.template,
        )

    def _add_title_slide(self, prs: Presentation, title: str) -> None:
        """Add the title slide with template styling."""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Set background
        self._set_slide_background(slide, self.style.bg_color)

        # Add accent bar at top
        self._add_accent_bar(slide)

        # Add title text
        left = Inches(1.0)
        top = Inches(2.5)
        width = Inches(11.0)
        height = Inches(2.0)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = tf.paragraphs[0]
        p.text = title
        p.font.name = self.style.title_font
        p.font.size = Pt(44)
        p.font.bold = self.style.title_bold
        p.font.color.rgb = self.style.title_color
        p.alignment = PP_ALIGN.CENTER

        # Add footer if template has one
        if self.style.footer_text:
            self._add_footer(slide, self.style.footer_text)

        # Add logo placeholder for Corporate template
        if self.template == SlideTemplate.CORPORATE:
            self._add_logo_placeholder(slide)

    def _add_content_slide(self, prs: Presentation, section: SlideSection) -> None:
        """Add a content slide with title and bullet points."""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Set background
        self._set_slide_background(slide, self.style.bg_color)

        # Add accent bar
        self._add_accent_bar(slide)

        # Add title
        left = Inches(0.8)
        top = Inches(0.5)
        width = Inches(11.5)
        height = Inches(1.2)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = section.title
        p.font.name = self.style.title_font
        p.font.size = self.style.title_size
        p.font.bold = self.style.title_bold
        p.font.color.rgb = self.style.title_color
        p.alignment = PP_ALIGN.LEFT

        # Add bullet points
        if section.bullets:
            body_left = Inches(1.0)
            body_top = Inches(2.0)
            body_width = Inches(11.0)
            body_height = Inches(4.8)
            body_box = slide.shapes.add_textbox(
                body_left, body_top, body_width, body_height
            )
            body_tf = body_box.text_frame
            body_tf.word_wrap = True

            for i, bullet in enumerate(section.bullets):
                if i == 0:
                    p = body_tf.paragraphs[0]
                else:
                    p = body_tf.add_paragraph()

                p.text = f"• {bullet}"
                p.font.name = self.style.body_font
                p.font.size = self.style.body_size
                p.font.color.rgb = self.style.body_color
                p.space_after = Pt(12)
                p.alignment = PP_ALIGN.LEFT

        # Add footer if template has one
        if self.style.footer_text:
            self._add_footer(slide, self.style.footer_text)

    def _set_slide_background(self, slide, color: RGBColor) -> None:
        """Set the slide background to a solid color."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_accent_bar(self, slide) -> None:
        """Add a colored accent bar at the top of the slide."""
        left = Inches(0)
        top = Inches(0)
        width = Inches(13.333)
        height = Inches(0.08)

        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            left, top, width, height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.style.accent_color
        shape.line.fill.background()

    def _add_footer(self, slide, text: str) -> None:
        """Add footer text at the bottom of the slide."""
        left = Inches(0.8)
        top = Inches(7.0)
        width = Inches(11.5)
        height = Inches(0.4)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame

        p = tf.paragraphs[0]
        p.text = text
        p.font.name = self.style.body_font
        p.font.size = Pt(10)
        p.font.color.rgb = self.style.accent_color
        p.alignment = PP_ALIGN.RIGHT

    def _add_logo_placeholder(self, slide) -> None:
        """Add a logo placeholder shape (Corporate template only)."""
        left = Inches(0.8)
        top = Inches(0.5)
        width = Inches(2.0)
        height = Inches(1.0)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame

        p = tf.paragraphs[0]
        p.text = "[Company Logo]"
        p.font.name = "Calibri"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0x99, 0xAA, 0xBB)
        p.font.italic = True
        p.alignment = PP_ALIGN.LEFT

    @staticmethod
    def _sanitize_filename(title: str) -> str:
        """Convert a title to a safe filename."""
        # Remove non-alphanumeric characters (keep spaces and hyphens)
        safe = re.sub(r"[^\w\s-]", "", title)
        # Replace whitespace with hyphens
        safe = re.sub(r"\s+", "-", safe.strip())
        # Limit length
        safe = safe[:80] if safe else "presentation"
        return f"{safe}.pptx"


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def generate_slides(
    summary: str,
    template: SlideTemplate = SlideTemplate.MODERN,
    title: Optional[str] = None,
) -> SlideResult:
    """
    Generate a PPTX presentation from a summary.

    This is the main entry point for slide generation. It parses the summary
    into sections, applies the chosen template styling, and returns the
    generated PPTX as bytes.

    Args:
        summary: The summary text to convert into slides.
        template: The visual template to apply (default: Modern).
        title: Optional presentation title override.

    Returns:
        SlideResult containing file bytes, filename, slide count, and template used.
    """
    generator = SlideGenerator(template=template)
    return generator.generate(summary, title=title)
