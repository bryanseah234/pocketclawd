"""
Draft artifact generation — convert markdown draft text from /draft into
a downloadable .docx (minutes/research/email) or .pptx (slides) file.

Used by the /draft command (R8 / Wave 11): the file is uploaded to S3 via
the data-gateway worker and the user receives a presigned URL.
"""
from __future__ import annotations

import io
import logging
import re
from datetime import datetime, timezone

logger = logging.getLogger(__name__)


def _sanitize_topic(topic: str, max_len: int = 60) -> str:
    """Make a topic string safe for use in a filename."""
    safe = re.sub(r"[^A-Za-z0-9_-]+", "-", topic.strip())
    safe = safe.strip("-")
    if not safe:
        safe = "draft"
    return safe[:max_len].lower()


def make_filename(doc_type: str, topic: str) -> str:
    """Build a deterministic filename for a given draft."""
    ext = {"slides": "pptx", "research": "pdf"}.get(doc_type, "docx")
    ts = datetime.now(timezone.utc).strftime("%Y%m%dT%H%M%SZ")
    safe = _sanitize_topic(topic)
    return f"{doc_type}-{safe}-{ts}.{ext}"


def render_docx(doc_type: str, topic: str, body_md: str) -> bytes:
    """
    Render a Word document from the LLM markdown output.

    The structure is intentionally simple: H1 = topic, then each
    `# heading` / `## heading` line in the markdown becomes a heading,
    each `- bullet` line becomes a list paragraph, and any other line is
    a regular paragraph. Empty lines are dropped.
    """
    try:
        from docx import Document
        from docx.shared import Pt
    except ImportError as exc:
        raise RuntimeError("python-docx is not installed") from exc

    doc = Document()
    style = doc.styles["Normal"]
    style.font.size = Pt(11)
    title_text = topic.strip() or doc_type.capitalize()
    doc.add_heading(title_text, level=0)
    doc.add_paragraph(
        f"Draft type: {doc_type} - generated {datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M UTC')}"
    )
    doc.add_paragraph("")

    for raw_line in body_md.splitlines():
        line = raw_line.rstrip()
        if not line.strip():
            continue
        if line.startswith("### "):
            doc.add_heading(line[4:].strip(), level=3)
        elif line.startswith("## "):
            doc.add_heading(line[3:].strip(), level=2)
        elif line.startswith("# "):
            doc.add_heading(line[2:].strip(), level=1)
        elif line.lstrip().startswith(("- ", "* ")):
            text = line.lstrip()[2:]
            doc.add_paragraph(text, style="List Bullet")
        else:
            doc.add_paragraph(line)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _split_slides(body_md: str) -> list[tuple[str, list[str]]]:
    """Parse the LLM slide outline into (title, bullets) tuples."""
    slides: list[tuple[str, list[str]]] = []
    current_title: str | None = None
    current_bullets: list[str] = []

    def flush() -> None:
        nonlocal current_title, current_bullets
        if current_title is not None:
            slides.append((current_title, current_bullets))
        current_title = None
        current_bullets = []

    for raw in body_md.splitlines():
        line = raw.rstrip()
        if not line.strip():
            continue
        # A new slide starts on # / ## / ### / Slide N: / numbered like "1." or "1)"
        m = re.match(r"^\s*(?:#{1,3}|Slide\s*\d+\s*[:.\-]|\d+\s*[\.\)])\s*(.+)$", line)
        if m and not line.lstrip().startswith(("-", "*")):
            flush()
            current_title = m.group(1).strip().lstrip(":-").strip()
            current_bullets = []
            continue
        # Bullet
        if line.lstrip().startswith(("- ", "* ")):
            if current_title is None:
                current_title = "Slide"
            current_bullets.append(line.lstrip()[2:].strip())
            continue
        # Plain line — append to bullets
        if current_title is None:
            current_title = line.strip()[:80]
            continue
        current_bullets.append(line.strip())

    flush()
    if not slides:
        slides = [(body_md.strip()[:80] or "Slide", [])]
    return slides



def render_pdf(topic: str, body_md: str) -> bytes:
    """
    Render a research report to PDF using reportlab.

    Simple two-pass: title heading, then body paragraphs. Markdown headings
    become bold/larger; bullets become indented lines. No images.
    """
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import cm
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.enums import TA_LEFT
    except ImportError as exc:
        raise RuntimeError("reportlab is not installed") from exc

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4,
                            leftMargin=2*cm, rightMargin=2*cm,
                            topMargin=2.5*cm, bottomMargin=2*cm)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("rtitle", parent=styles["Title"], fontSize=18, spaceAfter=12)
    h1_style    = ParagraphStyle("rh1",    parent=styles["Heading1"], fontSize=14, spaceAfter=6)
    h2_style    = ParagraphStyle("rh2",    parent=styles["Heading2"], fontSize=12, spaceAfter=4)
    body_style  = ParagraphStyle("rbody",  parent=styles["Normal"],   fontSize=10, spaceAfter=4, leading=14)
    bullet_style= ParagraphStyle("rbullet",parent=styles["Normal"],   fontSize=10, leftIndent=18, spaceAfter=3, leading=13)
    meta_style  = ParagraphStyle("rmeta",  parent=styles["Normal"],   fontSize=8,  textColor=(0.4,0.4,0.4), spaceAfter=12)

    story = [
        Paragraph(topic.strip() or "Research Report", title_style),
        Paragraph(
            f"Generated {datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M UTC')}",
            meta_style,
        ),
        Spacer(1, 0.3*cm),
    ]

    for raw_line in body_md.splitlines():
        line = raw_line.rstrip()
        if not line.strip():
            story.append(Spacer(1, 0.2*cm))
            continue
        # Escape XML special chars for reportlab
        safe = line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        if line.startswith("### "):
            story.append(Paragraph(safe[4:], h2_style))
        elif line.startswith("## "):
            story.append(Paragraph(safe[3:], h2_style))
        elif line.startswith("# "):
            story.append(Paragraph(safe[2:], h1_style))
        elif line.lstrip().startswith(("- ", "* ")):
            story.append(Paragraph(f"• {safe.lstrip()[2:]}", bullet_style))
        else:
            story.append(Paragraph(safe, body_style))

    doc.build(story)
    return buf.getvalue()


def render_pptx(topic: str, body_md: str) -> bytes:
    """Render a slide deck from the markdown outline."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError as exc:
        raise RuntimeError("python-pptx is not installed") from exc

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    title_layout = prs.slide_layouts[0]
    s0 = prs.slides.add_slide(title_layout)
    if s0.shapes.title is not None:
        s0.shapes.title.text = topic.strip() or "Draft"
    if len(s0.placeholders) >= 2:
        s0.placeholders[1].text = (
            f"Generated {datetime.now(timezone.utc).strftime('%Y-%m-%d')}"
        )

    bullet_layout = prs.slide_layouts[1]  # Title + Content
    for title, bullets in _split_slides(body_md):
        slide = prs.slides.add_slide(bullet_layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = title[:90]
        if len(slide.placeholders) >= 2:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            if not bullets:
                tf.paragraphs[0].text = ""
            else:
                tf.paragraphs[0].text = bullets[0]
                tf.paragraphs[0].font.size = Pt(20)
                for b in bullets[1:]:
                    p = tf.add_paragraph()
                    p.text = b
                    p.font.size = Pt(20)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def render_artifact(doc_type: str, topic: str, body_md: str) -> tuple[bytes, str]:
    """
    Build the artifact bytes + content-type for a given draft type.
    """
    if doc_type == "slides":
        try:
            data = render_pptx(topic, body_md)
            return data, "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        except Exception as exc:
            logger.error("PPTX render failed for topic=%s: %s", topic, exc)
            raise
    if doc_type == "research":
        try:
            data = render_pdf(topic, body_md)
            return data, "application/pdf"
        except Exception as exc:
            # reportlab not available → fall back to docx silently
            logger.warning("PDF render failed, falling back to docx: %s", exc)
    # Default: word doc
    data = render_docx(doc_type, topic, body_md)
    return data, "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
