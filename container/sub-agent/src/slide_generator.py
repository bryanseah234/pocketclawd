"""
Slide generator — creates PowerPoint presentations from structured content.
Templates: executive_summary, technical_deep_dive, project_proposal, data_report
"""
import io
import json
import logging
from dataclasses import dataclass, field
from enum import Enum

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False
    logger.warning("python-pptx not installed — slide generation disabled")


class SlideTemplate(Enum):
    EXECUTIVE_SUMMARY = "executive_summary"
    TECHNICAL_DEEP_DIVE = "technical_deep_dive"
    PROJECT_PROPOSAL = "project_proposal"
    DATA_REPORT = "data_report"


# Colour schemes per template: (title_bg_rgb, accent_rgb)
_SCHEMES = {
    SlideTemplate.EXECUTIVE_SUMMARY:   ((0, 32, 96), (212, 175, 55)),   # navy / gold
    SlideTemplate.TECHNICAL_DEEP_DIVE: ((45, 45, 45), (0, 188, 212)),   # dark grey / cyan
    SlideTemplate.PROJECT_PROPOSAL:    ((27, 94, 32), (255, 255, 255)), # forest green / white
    SlideTemplate.DATA_REPORT:         ((38, 38, 38), (255, 100, 0)),   # charcoal / orange
}


@dataclass
class SlideContent:
    title: str
    bullet_points: list[str] = field(default_factory=list)
    notes: str = ""


@dataclass
class PresentationRequest:
    title: str
    template: SlideTemplate
    slides: list[SlideContent]
    author: str = "Clawd"


def _rgb(triple: tuple) -> "RGBColor":
    from pptx.dml.color import RGBColor as RC
    return RC(*triple)


def _add_title_slide(prs: "Presentation", title: str, scheme: tuple, author: str) -> None:
    from pptx.util import Inches, Pt
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = _rgb(scheme[0])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = _rgb(scheme[1])
    # Author
    aBox = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.5))
    aBox.text_frame.paragraphs[0].text = author
    aBox.text_frame.paragraphs[0].font.size = Pt(14)
    aBox.text_frame.paragraphs[0].font.color.rgb = _rgb((200, 200, 200))


def _add_content_slide(
    prs: "Presentation", content: SlideContent, scheme: tuple
) -> None:
    from pptx.util import Inches, Pt
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = _rgb((250, 250, 250))
    # Title bar
    tBar = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(10), Inches(1.2))
    tBar.fill.solid()
    tBar.fill.fore_color.rgb = _rgb(scheme[0])
    tf = tBar.text_frame
    tf.paragraphs[0].text = content.title
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = _rgb(scheme[1])
    # Bullets
    if content.bullet_points:
        bBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5))
        btf = bBox.text_frame
        btf.word_wrap = True
        for i, bp in enumerate(content.bullet_points):
            p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            p.text = f"• {bp}"
            p.font.size = Pt(16)
            p.font.color.rgb = _rgb((30, 30, 30))
    if content.notes:
        slide.notes_slide.notes_text_frame.text = content.notes


def generate_presentation(request: PresentationRequest) -> bytes:
    """Generate a .pptx file and return it as bytes."""
    if not _PPTX_AVAILABLE:
        raise RuntimeError("python-pptx is not installed")
    from pptx import Presentation as PR
    prs = PR()
    prs.core_properties.author = request.author
    prs.core_properties.title = request.title
    scheme = _SCHEMES[request.template]
    _add_title_slide(prs, request.title, scheme, request.author)
    for slide_content in request.slides:
        _add_content_slide(prs, slide_content, scheme)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def parse_slide_request(user_message: str, llm_response: str) -> "PresentationRequest | None":
    """Parse LLM JSON response into a PresentationRequest."""
    try:
        data = json.loads(llm_response)
        template_str = data.get("template", "executive_summary")
        try:
            template = SlideTemplate(template_str)
        except ValueError:
            template = SlideTemplate.EXECUTIVE_SUMMARY
        slides = [
            SlideContent(
                title=s.get("title", ""),
                bullet_points=s.get("bullet_points", []),
                notes=s.get("notes", ""),
            )
            for s in data.get("slides", [])
        ]
        return PresentationRequest(
            title=data["title"],
            template=template,
            slides=slides,
            author=data.get("author", "Clawd"),
        )
    except (KeyError, json.JSONDecodeError, TypeError):
        return None
